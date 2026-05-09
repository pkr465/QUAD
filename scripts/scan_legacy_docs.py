"""Scan manually-authored docx/pptx for stale facts that need patching."""
from __future__ import annotations

import io
import sys
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from docx import Document
from pptx import Presentation


TARGETS_DOCX = [
    Path(r"C:\work\05\QUAD\docs\Design_Document_QUAD_Agent.docx"),
    Path(r"C:\work\05\QUAD\docs\QUAD_Platform_Design_v2.docx"),
    Path(r"C:\work\05\QUAD\docs\QUAD_Sample_App_Prompts.docx"),
    Path(r"C:\work\05\QUAD\docs\USAGE_GUIDE.docx"),
    Path(r"C:\work\05\QUAD\docs\PRD_Qualcomm_DevWorkflows_v3.docx"),
]
TARGETS_PPTX = [
    Path(r"C:\work\05\QUAD\docs\QUAD_Executive_Pitch.pptx"),
]

NEEDLES = ["933", "90 tests", " 90 ", "April 2026", "v0.3", "0.3.0", "Phase A", "mock-only"]


def scan_runs(par_runs_iter, label: str, hits: dict) -> None:
    text = "".join(r.text for r in par_runs_iter)
    for n in NEEDLES:
        if n in text:
            hits.setdefault(n, []).append(f"{label}: {text[:140]}")


def scan_docx(p: Path) -> None:
    print(f"\n=== {p.name} ===")
    if not p.exists():
        print("  (missing)")
        return
    d = Document(str(p))
    hits: dict[str, list[str]] = {}
    for i, para in enumerate(d.paragraphs):
        scan_runs(para.runs, f"para[{i}]", hits)
    for ti, t in enumerate(d.tables):
        for ri, row in enumerate(t.rows):
            for ci, cell in enumerate(row.cells):
                for para in cell.paragraphs:
                    scan_runs(para.runs, f"table[{ti}].cell[{ri},{ci}]", hits)
    if not hits:
        print("  no stale needles found")
    for n, lines in hits.items():
        print(f"  [{n}]: {len(lines)} occurrences")
        for line in lines[:3]:
            print(f"     {line}")


def scan_pptx(p: Path) -> None:
    print(f"\n=== {p.name} ===")
    if not p.exists():
        print("  (missing)")
        return
    pr = Presentation(str(p))
    hits: dict[str, list[str]] = {}
    for si, slide in enumerate(pr.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    text = "".join(r.text for r in para.runs)
                    for n in NEEDLES:
                        if n in text:
                            hits.setdefault(n, []).append(
                                f"slide {si} para {pi}: {text[:140]}"
                            )
    if not hits:
        print("  no stale needles found")
    for n, lines in hits.items():
        print(f"  [{n}]: {len(lines)} occurrences")
        for line in lines[:3]:
            print(f"     {line}")


for f in TARGETS_DOCX:
    scan_docx(f)
for f in TARGETS_PPTX:
    scan_pptx(f)
