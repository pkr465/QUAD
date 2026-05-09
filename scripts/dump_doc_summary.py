"""Print a survey of QUAD/QUAD-Client design docs and pptx contents (UTF-8 safe)."""
from __future__ import annotations

import io
import sys
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from docx import Document
from pptx import Presentation


def survey_docx(path: Path) -> None:
    print(f"\n=== DOCX: {path} ===")
    if not path.exists():
        print("  (missing)")
        return
    d = Document(str(path))
    print(f"  paragraphs={len(d.paragraphs)}  tables={len(d.tables)}")
    for p in d.paragraphs[:18]:
        if p.text.strip():
            print(f"   {p.text[:160]!r}")


def survey_pptx(path: Path) -> None:
    print(f"\n=== PPTX: {path} ===")
    if not path.exists():
        print("  (missing)")
        return
    p = Presentation(str(path))
    print(f"  slides={len(p.slides)}")
    for i, slide in enumerate(p.slides, start=1):
        title = ""
        body_bits = []
        for s in slide.shapes:
            if s.has_text_frame and s.text_frame.text.strip():
                t = s.text_frame.text.strip()
                if not title:
                    title = t.splitlines()[0]
                body_bits.append(t[:120].replace("\n", " | "))
        print(f"  Slide {i}: {title[:80]}")
        for b in body_bits[1:4]:
            print(f"     - {b}")


roots = [
    Path(r"C:\work\05\QUAD\docs"),
    Path(r"C:\work\05\QUAD-Client\docs"),
]
for root in roots:
    print(f"\n###### ROOT: {root} ######")
    for f in sorted(root.glob("*.docx")):
        survey_docx(f)
    for f in sorted(root.glob("*.pptx")):
        survey_pptx(f)
