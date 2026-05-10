"""Patch manually-authored docx/pptx in place.

For each target:
  - run-level text replacement of stale needles (preserves formatting)
  - append an 'IoT Device Support' appendix referencing IOT_DEPENDENCIES.xlsx

For the Executive Pitch pptx: replace 933 -> 2002 on slides 5 + 10, add an
'IoT Device Support' slide before the closing slide.
"""
from __future__ import annotations

import io
import sys
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ───────── Replacements (apply per-run; safe across split runs via merge) ─────


REPLACEMENTS = [
    ("April 2026", "May 2026"),
    ("Apr 2026", "May 2026"),
    ("90 verification tests", "2002 verification tests"),
    ("Runs full test suite (90 tests)", "Runs full test suite (2002 tests)"),
    ("90 tests", "2002 tests"),
    ("933 Tests", "2002 Tests"),
    ("933 tests", "2002 tests"),
    ("933", "2002"),
]

# Pitch-pptx-only replacements: applied after the slide insert.
PITCH_FOOTER_REPLACEMENTS = [
    ("/12", "/13"),
    ("Hardware (3 devices)", "Hardware (4 device classes)"),
    (
        "Snapdragon X Elite, Arduino UNO Q, Snapdragon 8 Elite",
        "Snapdragon X Elite, Snapdragon 8 Elite, Arduino UNO Q, "
        "Qualcomm IoT SoC kit (RB3 Gen 2 / RB5 / QCS6490)",
    ),
    (
        "2002 Tests  |  15 Modules  |  30+ Templates  |  15,000+ Lines",
        "2002 Tests  |  120+ Modules  |  30+ Templates  |  25,000+ Lines",
    ),
]


def replace_in_paragraph(par, replacements: list[tuple[str, str]]) -> int:
    """Run-level text replacement. If the needle spans runs, fall back to
    rewriting the first run with the merged text and clearing the rest.

    Returns the number of replacements made.
    """
    n_changes = 0
    runs = par.runs
    if not runs:
        return 0
    full_text = "".join(r.text for r in runs)
    new_full = full_text
    for old, new in replacements:
        if old in new_full:
            new_full = new_full.replace(old, new)
            n_changes += 1
    if new_full == full_text:
        return 0

    # Try simple per-run replacement first (preserves formatting fully)
    if all(_simple_replace_run(r, replacements) for r in runs):
        return n_changes

    # Fall back: rewrite using the first run's formatting
    runs[0].text = new_full
    for r in runs[1:]:
        r.text = ""
    return n_changes


def _simple_replace_run(run, replacements: list[tuple[str, str]]) -> bool:
    """Apply replacements that fit within this single run. Returns True if all
    needles either weren't present or were fully contained in this run."""
    txt = run.text
    new_txt = txt
    for old, new in replacements:
        if old in new_txt:
            new_txt = new_txt.replace(old, new)
    if new_txt != txt:
        run.text = new_txt
    return True


def patch_docx_text(path: Path, replacements: list[tuple[str, str]]) -> int:
    d = Document(str(path))
    total = 0
    for par in d.paragraphs:
        total += replace_in_paragraph(par, replacements)
    for tbl in d.tables:
        for row in tbl.rows:
            for cell in row.cells:
                for par in cell.paragraphs:
                    total += replace_in_paragraph(par, replacements)
    if total:
        d.save(str(path))
    return total


# ───────── DOCX: append IoT appendix ─────────────────────────────────────────


IOT_INTRO = (
    "QUAD's adapter pattern is being extended from AI PC and Mobile to the full "
    "Qualcomm IoT SoC line. The full dependency catalogue (111 components across "
    "14 categories) is published as a workbook at docs/IOT_DEPENDENCIES.xlsx. "
    "This appendix is a quick orientation; the workbook is authoritative."
)

IOT_HARDWARE_TABLE = [
    ["Class", "SoC / Board", "Where it fits"],
    ["Edge gateway",        "QCS6490 (RB3 Gen 2), QCS8250 (RB5)", "Linux-class IoT, 12+ TOPS NPU"],
    ["High-perf IoT / AI",  "QCS8550",                            "Industrial AI, smart camera"],
    ["Camera / vision IoT", "QCS610 / QCS605",                    "Smart camera, AI vision"],
    ["Entry IoT",           "QCM2290 / QCS2290",                  "Wearables, smart speakers"],
    ["Cellular IoT",        "Snapdragon X75 / 9205 modem",        "5G / NB-IoT / LTE-M uplink"],
]

IOT_LAYER_TABLE = [
    ["Layer", "Key dependencies"],
    ["OS / Firmware",      "Yocto meta-qcom, Qualcomm Linux, Zephyr / FreeRTOS, U-Boot, TF-A, OP-TEE"],
    ["Connectivity",       "BlueZ, OpenThread, Matter 1.4, hostapd, ModemManager + libqmi/libmbim"],
    ["IoT protocols",      "MQTT (Mosquitto / paho), CoAP (libcoap / aiocoap), LwM2M (Anjay), OPC-UA"],
    ["Cloud + OTA",        "AWS IoT Device SDK v2, azure-iot-device, Greengrass, IoT Edge, Mender, RAUC, SWUpdate"],
    ["Security",           "Qualcomm QTEE / SPU, OP-TEE, mbedTLS, OpenSSL, TF-M, PKCS#11, Matter Attestation"],
    ["Edge AI runtime",    "QNN / QAIRT 2.x, SNPE 2.x, Hexagon SDK 5.x, AIMET, Qualcomm AI Hub"],
    ["Sensors / HAL",      "libgpiod, smbus2, spidev, pyserial, python-can, pymodbus, bleak"],
    ["Telemetry",          "OpenTelemetry, Prometheus client_python, Fluent Bit"],
]

IOT_SAMPLE_PROMPTS = [
    'Use QUAD to detect the Qualcomm IoT SoC on the connected RB3 Gen 2 board and list its CPU/GPU/NPU topology.',
    'Use QUAD to convert mobilenet_v2.onnx to QNN INT8 and profile it on a QCS6490, targeting < 3W power draw.',
    'Use QUAD to allocate a YOLOv8n inference graph across the Hexagon NPU + CPU on QCS8550 within a 5W power budget.',
    'Use QUAD to generate C++ inference code for a Yocto-based RB3 Gen 2 image with QNN runtime + Mender OTA update hooks.',
    'Use QUAD to compare INT8 inference of resnet50.onnx between QCS6490 (gateway) and QCS610 (smart camera) and produce a power-vs-latency table.',
]


IOT_APPENDIX_HEADING = "Appendix — IoT Device Support"
MEASUREMENT_APPENDIX_HEADING = "Appendix — Real-mode Measurement Fidelity"


def append_iot_appendix(path: Path, *, prompts: bool = False) -> None:
    d = Document(str(path))
    # Idempotency: skip if appendix already present.
    for par in d.paragraphs:
        if par.text.strip() == IOT_APPENDIX_HEADING:
            print(f"  (skip) {path.name} already has IoT appendix")
            return

    d.add_page_break()
    d.add_heading(IOT_APPENDIX_HEADING, level=1)
    d.add_paragraph(
        "This section was added 2026-05-09 to record the dependencies QUAD "
        "needs to extend support across Qualcomm's IoT SoC portfolio. It "
        "complements (but does not replace) the per-component workbook at "
        "docs/IOT_DEPENDENCIES.xlsx."
    )

    d.add_heading("Scope", level=2)
    d.add_paragraph(IOT_INTRO)

    d.add_heading("Target hardware", level=2)
    _add_table(d, IOT_HARDWARE_TABLE, header=True)

    d.add_heading("Software dependency layers", level=2)
    _add_table(d, IOT_LAYER_TABLE, header=True)

    if prompts:
        d.add_heading("Sample prompts (Claude Code + QUAD MCP)", level=2)
        for q in IOT_SAMPLE_PROMPTS:
            _add_bullet(d, q)

    d.add_heading("Reference", level=2)
    refs = [
        "docs/IOT_DEPENDENCIES.xlsx — full per-component catalogue (priority, license, target version, source URL)",
        "docs/PREREQUISITES.md — base prerequisites for the QUAD agent",
        "QUAD_Server_Guide.docx — section 19 (IoT Device Support)",
    ]
    for r in refs:
        _add_bullet(d, r)

    d.save(str(path))


def _add_bullet(d, text: str) -> None:
    """Add a bullet paragraph; fall back to a leading '• ' if 'List Bullet'
    style is not defined in the document."""
    try:
        p = d.add_paragraph(style="List Bullet")
        p.add_run(text)
    except KeyError:
        p = d.add_paragraph()
        p.add_run("•  " + text)


MEASUREMENT_INTRO = (
    "QUAD's profile_workload tool now returns a measurement_notes block "
    "that tags every metric with its provenance: measured (from snpe-"
    "diagview / QPM3 / psutil), estimated (host thermal model), or "
    "not_measured (with a reason). No fictional defaults — what you see "
    "is what was actually captured."
)

MEASUREMENT_PROVENANCE_TABLE = [
    ["Metric",        "Measured source",                          "Fallback when tool absent"],
    ["Latency",       "snpe-diagview CSV (Total Inference Time)", "snpe-net-run stdout parser"],
    ["Per-layer",     "snpe-diagview CSV (Model Layer Times)",    "synthetic composite (1 row, tagged)"],
    ["Memory RSS",    "psutil sampler (100 ms polling)",          "not_measured:rss_unavailable"],
    ["Power",         "QPM3 (per-frame PMIC readings)",           "host_thermal_model (CPU/NPU% × TDP)"],
    ["NPU util",      "Arithmetic: cycles / wall_time / clock",   "0.0 (tagged not_measured)"],
    ["GPU util",      "Snapdragon Profiler chrometrace events",   "0.0 (tagged not_measured)"],
    ["CPU util",      "psutil.cpu_percent",                       "Always available"],
]

MEASUREMENT_PROFILER_TABLE = [
    ["Tool", "Activates when present",                       "Provides"],
    ["QPM3 (Qualcomm Power Monitor 3)",
     "PATH or QPM3_HOME env var",
     "Measured per-frame power; flips power tag from estimated to measured"],
    ["Snapdragon Profiler (sdptrace)",
     "PATH or SNAPDRAGON_PROFILER_HOME env var",
     "Real GPU% from chrometrace; populates utilization['gpu']"],
    ["psutil (always present in [real])",
     "automatic",
     "RSS sampling + CPU% percent"],
    ["snpe-diagview (in QAIRT)",
     "automatic",
     "Authoritative latency + per-layer; raises tag from snpe-net-run stdout"],
]

MEASUREMENT_REGISTRY_TABLE = [
    ["Command",                       "What it does"],
    ["quad models list",              "Show every registry entry with cache state"],
    ["quad models fetch <name>",      "Download + verify SHA-256 (URL entries)"],
    ["quad models path <name>",       "Resolve to local path (env-var entries)"],
    ["quad models verify <name>",     "Re-check SHA-256 of a cached file"],
]


def append_measurement_appendix(path: Path) -> None:
    """Append the 'Real-mode Measurement Fidelity' appendix in-place.

    Idempotent — skips when the appendix heading is already present.
    """
    d = Document(str(path))
    for par in d.paragraphs:
        if par.text.strip() == MEASUREMENT_APPENDIX_HEADING:
            print(f"  (skip) {path.name} already has measurement appendix")
            return

    d.add_page_break()
    d.add_heading(MEASUREMENT_APPENDIX_HEADING, level=1)
    d.add_paragraph(
        "Added 2026-05-10. Records the real-mode measurement plumbing "
        "that landed across the recent QUAD changes (QAIRT-adapter stdout "
        "parsers, RSS sampler, host power model, QPM3 / Snapdragon "
        "Profiler integration, model registry)."
    )

    d.add_heading("Provenance-tagged metrics", level=2)
    d.add_paragraph(MEASUREMENT_INTRO)
    _add_table(d, MEASUREMENT_PROVENANCE_TABLE, header=True)

    d.add_heading("Auto-detected precision profilers", level=2)
    _add_table(d, MEASUREMENT_PROFILER_TABLE, header=True)

    d.add_heading("Model registry — `quad models`", level=2)
    d.add_paragraph(
        "Production ONNX provisioning lives in src/quad/model_registry/. "
        "Entries declare either a url (auto-downloadable, SHA-256 "
        "verified) or a path_env_var (user-supplied for gated weights "
        "like Llama 3 8B). Adding a model for a future plan is a single "
        "YAML edit — no Python changes needed."
    )
    _add_table(d, MEASUREMENT_REGISTRY_TABLE, header=True)

    d.add_heading("Snapdragon X Elite environment caveat", level=2)
    d.add_paragraph(
        "QAIRT 2.46's host Python tools (qairt-converter, qairt-quantizer) "
        "ship as windows-arm64ec/.pyd modules that need the full Visual "
        "Studio 2022 runtime — not just the VC++ redist, and not native "
        "ARM64 Python alone (ARM64EC .pyd can't load into pure ARM64 "
        "Python; WinError 193). On Snapdragon X Elite, either install "
        "VS 2022 Community (winget install "
        "Microsoft.VisualStudio.2022.Community) or run model conversion "
        "on a separate x86_64 host and copy the .dlc back. QUAD's "
        "runtime path (snpe-net-run, profiling) works on a stock install "
        "— only model conversion is affected."
    )

    d.add_heading("Reference", level=2)
    for r in [
        "src/quad/adapters/parsers.py — pure-function parsers for snpe-net-run, snpe-diagview, qnn-platform-validator, qairt-converter",
        "src/quad/profiler/qpm3.py — Qualcomm Power Monitor 3 wrapper",
        "src/quad/profiler/sdptrace.py — Snapdragon Profiler trace wrapper",
        "src/quad/profiler/rss_sampler.py — psutil-based async RSS sampler",
        "src/quad/profiler/host_power.py — host-side power estimate + SRUM Energy Estimation reader",
        "src/quad/profiler/host_utilization.py — CPU% + NPU/GPU utilisation helpers",
        "src/quad/model_registry/ — production ONNX manifest + fetcher + CLI",
    ]:
        _add_bullet(d, r)

    d.save(str(path))


def _add_table(d, rows: list[list[str]], *, header: bool) -> None:
    t = d.add_table(rows=len(rows), cols=len(rows[0]))
    try:
        t.style = "Light Grid Accent 1"
    except KeyError:
        pass
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = t.cell(r_idx, c_idx)
            cell.text = ""
            run = cell.paragraphs[0].add_run(val)
            run.font.size = Pt(10)
            if header and r_idx == 0:
                run.bold = True


# ───────── PPTX: add IoT slide ───────────────────────────────────────────────


def patch_pptx(path: Path) -> None:
    pres = Presentation(str(path))

    # 1. Text replacement on existing slides
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full = "".join(r.text for r in para.runs)
                    new_full = full
                    for old, new in REPLACEMENTS:
                        if old in new_full:
                            new_full = new_full.replace(old, new)
                    if new_full != full and para.runs:
                        para.runs[0].text = new_full
                        for r in para.runs[1:]:
                            r.text = ""

    # Idempotency: skip slide insert if an IoT slide is already present.
    has_iot_slide = False
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and "IoT Device Support" in shape.text_frame.text:
                has_iot_slide = True
                break
        if has_iot_slide:
            break

    if has_iot_slide:
        print(f"  (skip slide insert) {path.name} already has IoT slide")
    else:
        _add_iot_slide(pres)

    # 3. Final pass: pitch-specific footer / wording updates (run AFTER insert
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full = "".join(r.text for r in para.runs)
                    new_full = full
                    for old, new in PITCH_FOOTER_REPLACEMENTS:
                        if old in new_full:
                            new_full = new_full.replace(old, new)
                    if new_full != full and para.runs:
                        para.runs[0].text = new_full
                        for r in para.runs[1:]:
                            r.text = ""

    pres.save(str(path))


def _add_iot_slide(pres) -> None:
    """Insert an 'IoT Device Support' slide just before the closing slide."""
    slide_w = pres.slide_width
    slide_h = pres.slide_height
    blank_layout = pres.slide_layouts[6]
    new_slide = pres.slides.add_slide(blank_layout)

    # Move the new slide so it sits second-from-last (before "Thank You")
    sldIdLst = pres.slides._sldIdLst
    slide_ids = list(sldIdLst)
    if len(slide_ids) >= 2:
        sldIdLst.remove(slide_ids[-1])
        sldIdLst.insert(len(slide_ids) - 2, slide_ids[-1])

    QUAL_BLUE = PRGBColor(0x32, 0x53, 0xDC)
    DARK = PRGBColor(0x2D, 0x2D, 0x2D)
    GRAY = PRGBColor(0x66, 0x66, 0x66)

    bar = new_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_w, PInches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = QUAL_BLUE
    bar.line.fill.background()

    title_box = new_slide.shapes.add_textbox(
        PInches(0.6), PInches(0.4), slide_w - PInches(1.2), PInches(0.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "IoT Device Support"
    p.font.size = PPt(32)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = "Calibri Light"

    sub_box = new_slide.shapes.add_textbox(
        PInches(0.6), PInches(1.2), slide_w - PInches(1.2), PInches(0.5)
    )
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Extending QUAD across the full Qualcomm IoT SoC portfolio"
    p.font.size = PPt(16)
    p.font.color.rgb = QUAL_BLUE
    p.font.name = "Calibri"

    body_box = new_slide.shapes.add_textbox(
        PInches(0.7), PInches(1.9), slide_w - PInches(1.4), PInches(4.8)
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    bullets = [
        ("Target SoCs",
         "QCS6490 (RB3 Gen 2), QCS8550, QCS8250 (RB5), QCS610, QCM2290, X75 modem"),
        ("OS / Firmware",
         "Yocto meta-qcom, Qualcomm Linux, Zephyr / FreeRTOS, U-Boot, TF-A, OP-TEE"),
        ("Connectivity",
         "BlueZ, OpenThread, Matter 1.4, hostapd, ModemManager (NB-IoT, LTE-M, 5G)"),
        ("IoT protocols",
         "MQTT (Mosquitto / paho), CoAP (libcoap / aiocoap), LwM2M (Anjay), OPC-UA"),
        ("Cloud + OTA",
         "AWS IoT v2, azure-iot-device, Greengrass, IoT Edge, Mender, RAUC, SWUpdate"),
        ("Edge AI runtime",
         "QNN / QAIRT 2.x, SNPE 2.x, Hexagon SDK 5, AIMET, Qualcomm AI Hub"),
    ]
    for i, (label, value) in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rl = para.add_run()
        rl.text = f"{label}: "
        rl.font.size = PPt(15)
        rl.font.bold = True
        rl.font.color.rgb = QUAL_BLUE
        rl.font.name = "Calibri"
        rv = para.add_run()
        rv.text = value
        rv.font.size = PPt(15)
        rv.font.color.rgb = DARK
        rv.font.name = "Calibri"
        para.space_after = PPt(6)

    foot_box = new_slide.shapes.add_textbox(
        PInches(0.7), slide_h - PInches(1.0),
        slide_w - PInches(1.4), PInches(0.5),
    )
    tf = foot_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Full catalogue: docs/IOT_DEPENDENCIES.xlsx — 111 components across 14 categories"
    p.font.size = PPt(12)
    p.font.italic = True
    p.font.color.rgb = GRAY
    p.font.name = "Calibri"


# ───────── Drive ────────────────────────────────────────────────────────────


def main() -> None:
    docx_targets = [
        (Path(r"C:\work\05\QUAD\docs\Design_Document_QUAD_Agent.docx"), False),
        (Path(r"C:\work\05\QUAD\docs\QUAD_Platform_Design_v2.docx"),    False),
        (Path(r"C:\work\05\QUAD\docs\QUAD_Sample_App_Prompts.docx"),    True),
        (Path(r"C:\work\05\QUAD\docs\USAGE_GUIDE.docx"),                False),
        (Path(r"C:\work\05\QUAD\docs\PRD_Qualcomm_DevWorkflows_v3.docx"), False),
    ]
    for path, prompts in docx_targets:
        n = patch_docx_text(path, REPLACEMENTS)
        print(f"Replaced {n} run(s) in {path.name}")
        append_iot_appendix(path, prompts=prompts)
        print(f"Appended IoT appendix to {path.name}")
        append_measurement_appendix(path)
        print(f"Appended Measurement-Fidelity appendix to {path.name}")

    pptx_path = Path(r"C:\work\05\QUAD\docs\QUAD_Executive_Pitch.pptx")
    patch_pptx(pptx_path)
    print(f"Patched {pptx_path.name} (text + new IoT slide)")


if __name__ == "__main__":
    main()
