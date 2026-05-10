"""
Generate QUAD User Journey PowerPoint presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os

# Constants
QUALCOMM_BLUE = RGBColor(0x32, 0x53, 0xDC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
GRAY_TEXT = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)

TITLE_FONT = "Calibri Light"
BODY_FONT = "Calibri"
TITLE_SIZE = Pt(32)
BODY_SIZE = Pt(18)
SUBTITLE_SIZE = Pt(14)

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

OUTPUT_PATH = str(Path(__file__).resolve().parent / "QUAD_User_Journey.pptx")


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout for full control
    blank_layout = prs.slide_layouts[6]  # Blank

    # Slide 1: Title Slide
    create_title_slide(prs, blank_layout)

    # Slide 2: Discovery
    create_content_slide(
        prs, blank_layout,
        title="Step 1: Discover QUAD",
        bullets=[
            "Developer finds QUAD via: GitHub, Qualcomm DevRel, Claude Code marketplace",
            'First impression: README with clear value proposition',
            '"One platform for all Qualcomm silicon"',
            "Unified SDK for NPU, GPU, DSP across mobile, PC, IoT, automotive",
        ],
        slide_num=2,
    )

    # Slide 3: Installation
    create_content_slide(
        prs, blank_layout,
        title="Step 2: Install (One Command)",
        bullets=[
            "./install.sh \u2014 installs everything in 3 minutes",
            "What happens: Python venv, QUAD package, SDK download, VS Code setup, tests verified",
            "Result: 2002 tests pass, Claude Code auto-detects MCP server",
            "No manual SDK download needed",
        ],
        slide_num=3,
        code_snippet="$ ./install.sh\n\u2713 Python venv created\n\u2713 QUAD installed\n\u2713 QNN SDK downloaded\n\u2713 2002 tests passed",
    )

    # Slide 4: Configuration
    create_content_slide(
        prs, blank_layout,
        title="Step 3: Configure",
        bullets=[
            "quad configure \u2014 interactive wizard asks 5 questions",
            "Outputs: quad.toml (paths) + .env (secrets)",
            "Or: mock mode works with ZERO configuration",
            "Automatic detection of installed SDKs and devices",
        ],
        slide_num=4,
    )

    # Slide 5: First Inference
    create_content_slide(
        prs, blank_layout,
        title='Step 4: First Inference \u2014 < 5 Minutes',
        bullets=[
            "quad quickstart \u2014 guided wizard:",
            "  1. Detects hardware (Snapdragon X Elite NPU @ 45 TOPS)",
            "  2. Compiles sample model (ONNX \u2192 QNN)",
            "  3. Profiles (latency: 5ms, power: 2W)",
            "  4. Generates inference code (C++/Python/Kotlin)",
        ],
        slide_num=5,
        footer_note="Time from install to running inference: < 5 minutes",
    )

    # Slide 6: Development Workflow
    create_content_slide(
        prs, blank_layout,
        title="Step 5: Build Your Model",
        bullets=[
            "Convert:  quad compile model.onnx --targets all",
            "Optimize: quad optimize --quantization int8 --power-budget 5W",
            "Profile:  quad profile --level deep (roofline + power + memory)",
            "Generate: quad generate --platform android --language kotlin",
        ],
        slide_num=6,
    )

    # Slide 7: AI Agent (MCP)
    create_content_slide(
        prs, blank_layout,
        title="Step 6: Talk to QUAD",
        bullets=[
            "Natural language in Claude Code:",
            '  "Convert my model to INT8 and deploy to the phone"',
            '  "Profile power on NPU vs GPU \u2014 use QPM3 when available"',
            '  "List models, fetch the latest mobilenetv2 from the registry"',
            "5 MCP tools auto-registered \u2014 zero config in Claude Code",
            "Every returned metric carries a provenance tag (measured/estimated/not_measured)",
        ],
        slide_num=7,
    )

    # Slide 8: Deploy to Device
    create_content_slide(
        prs, blank_layout,
        title="Step 7: Deploy",
        bullets=[
            "./deploy.sh model.onnx --runtime dsp --quantize int8",
            "  Step 1: Convert \u2192 DLC",
            "  Step 2: SCP to device",
            "  Step 3: snpe-net-run on target",
            "  Step 4: Retrieve results",
        ],
        slide_num=8,
        footer_note="Supports: Windows (local), Linux (SSH), Android (ADB), Qualcomm IoT (RB3 Gen 2 / RB5)",
    )

    # Slide 9: Production Serving
    create_content_slide(
        prs, blank_layout,
        title="Step 8: Serve in Production",
        bullets=[
            "QUAD Serve: ModelServer with dynamic batching",
            "quad serve --models ./repo/ --port 8080",
            "Health checks, metrics (Prometheus), model hot-swap",
            "Power-budget-aware scheduling",
        ],
        slide_num=9,
    )

    # Slide 10: Scaling Up
    create_content_slide(
        prs, blank_layout,
        title="Step 9: Scale & Optimize",
        bullets=[
            "PSNPE: parallel inference across NPU cores",
            "Custom kernels: @quad.kernel Python DSL \u2192 Hexagon",
            "Multi-model pipelines (detection + classification + segmentation)",
            "IoT fleet rollout: QCS6490 / QCS8550 / RB3 Gen 2 / RB5 \u2014 see docs/IOT_DEPENDENCIES.xlsx",
            "OTA via Mender / RAUC; MQTT + CoAP + LwM2M telemetry to AWS / Azure IoT",
        ],
        slide_num=10,
    )

    # Slide 11: Community
    create_content_slide(
        prs, blank_layout,
        title="Step 10: Join the Community",
        bullets=[
            "QUAD Academy: free courses",
            'Certification: "Qualcomm AI Developer Certified"',
            "Hackathons, Discord, StackOverflow tag",
            "Open-source samples (50+ on GitHub)",
        ],
        slide_num=11,
    )

    # Slide 12: Summary
    create_summary_slide(prs, blank_layout)

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")


def add_slide_number(slide, num, total=12):
    """Add slide number in footer."""
    left = SLIDE_WIDTH - Inches(1.2)
    top = SLIDE_HEIGHT - Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(1.0), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = BODY_FONT
    p.alignment = PP_ALIGN.RIGHT


def add_accent_bar(slide):
    """Add a thin blue accent bar at the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = QUALCOMM_BLUE
    shape.line.fill.background()


def create_title_slide(prs, layout):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(layout)

    # Blue accent bar at top
    add_accent_bar(slide)

    # Qualcomm logo placeholder (top-right)
    left = SLIDE_WIDTH - Inches(3.5)
    top = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(3.0), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "QUALCOMM"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = QUALCOMM_BLUE
    p.font.name = TITLE_FONT
    p.alignment = PP_ALIGN.RIGHT

    # Main title
    left = Inches(1.5)
    top = Inches(2.2)
    txBox = slide.shapes.add_textbox(left, top, Inches(10), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "QUAD User Journey"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.font.name = TITLE_FONT
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    top = Inches(3.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "From Zero to NPU Inference in 5 Minutes"
    p.font.size = Pt(24)
    p.font.color.rgb = QUALCOMM_BLUE
    p.font.name = BODY_FONT
    p.alignment = PP_ALIGN.LEFT

    # Date
    top = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, Inches(4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "May 2026"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = BODY_FONT
    p.alignment = PP_ALIGN.LEFT

    # Decorative blue rectangle (bottom-left)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), SLIDE_HEIGHT - Inches(1.0),
        Inches(4), Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = QUALCOMM_BLUE
    shape.line.fill.background()

    add_slide_number(slide, 1)


def create_content_slide(prs, layout, title, bullets, slide_num,
                         code_snippet=None, footer_note=None):
    """Create a standard content slide with title and bullets."""
    slide = prs.slides.add_slide(layout)

    # Accent bar
    add_accent_bar(slide)

    # Title
    left = Inches(0.8)
    top = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(11), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = TITLE_SIZE
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.font.name = TITLE_FONT
    p.alignment = PP_ALIGN.LEFT

    # Blue underline below title
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, Inches(1.35),
        Inches(2.5), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = QUALCOMM_BLUE
    shape.line.fill.background()

    # Bullets
    bullet_left = Inches(1.0)
    bullet_top = Inches(1.8)
    bullet_width = Inches(7.0) if code_snippet else Inches(11.0)
    bullet_height = Inches(4.5)

    txBox = slide.shapes.add_textbox(bullet_left, bullet_top, bullet_width, bullet_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Indented items
        if bullet.startswith("  "):
            p.text = bullet.strip()
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY_TEXT
            p.level = 1
            p.space_before = Pt(4)
        else:
            p.text = bullet
            p.font.size = BODY_SIZE
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(10)

        p.font.name = BODY_FONT
        p.space_after = Pt(4)

    # Code snippet (right side)
    if code_snippet:
        code_left = Inches(8.5)
        code_top = Inches(1.8)
        code_width = Inches(4.2)
        code_height = Inches(3.5)

        # Background rectangle for code
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            code_left, code_top,
            code_width, code_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
        bg_shape.line.fill.background()

        # Code text
        code_box = slide.shapes.add_textbox(
            code_left + Inches(0.2), code_top + Inches(0.2),
            code_width - Inches(0.4), code_height - Inches(0.4)
        )
        ctf = code_box.text_frame
        ctf.word_wrap = True
        for i, line in enumerate(code_snippet.split("\n")):
            if i == 0:
                cp = ctf.paragraphs[0]
            else:
                cp = ctf.add_paragraph()
            cp.text = line
            cp.font.size = Pt(12)
            cp.font.name = "Consolas"
            if line.startswith("\u2713"):
                cp.font.color.rgb = RGBColor(0x50, 0xFA, 0x7B)
            elif line.startswith("$"):
                cp.font.color.rgb = RGBColor(0xBD, 0x93, 0xF9)
            else:
                cp.font.color.rgb = RGBColor(0xF8, 0xF8, 0xF2)

    # Footer note
    if footer_note:
        note_left = Inches(1.0)
        note_top = SLIDE_HEIGHT - Inches(1.2)
        txBox = slide.shapes.add_textbox(note_left, note_top, Inches(10), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = footer_note
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = QUALCOMM_BLUE
        p.font.name = BODY_FONT
        p.alignment = PP_ALIGN.LEFT

    add_slide_number(slide, slide_num)


def create_summary_slide(prs, layout):
    """Slide 12: Summary with flow diagram."""
    slide = prs.slides.add_slide(layout)

    # Accent bar
    add_accent_bar(slide)

    # Title
    left = Inches(0.8)
    top = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(11), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "The Complete Journey"
    p.font.size = TITLE_SIZE
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.font.name = TITLE_FONT

    # Blue underline
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, Inches(1.35),
        Inches(2.5), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = QUALCOMM_BLUE
    shape.line.fill.background()

    # Flow diagram: 10 steps in two rows
    steps = [
        "Discover", "Install", "Configure", "First\nInference", "Develop",
        "AI Agent", "Deploy", "Serve", "Scale", "Community"
    ]

    # Row 1: steps 1-5
    start_x = Inches(0.6)
    y1 = Inches(2.2)
    y2 = Inches(4.2)
    box_w = Inches(2.0)
    box_h = Inches(1.2)
    gap = Inches(0.4)

    for row_idx, y in enumerate([y1, y2]):
        row_steps = steps[row_idx * 5:(row_idx + 1) * 5]
        for i, step_name in enumerate(row_steps):
            step_num = row_idx * 5 + i + 1
            x = start_x + i * (box_w + gap)

            # Rounded rectangle
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, box_w, box_h
            )
            box.fill.solid()
            if step_num == 4:  # Highlight "First Inference" as the aha moment
                box.fill.fore_color.rgb = QUALCOMM_BLUE
                text_color = WHITE
            else:
                box.fill.fore_color.rgb = RGBColor(0xE8, 0xEB, 0xF7)
                text_color = DARK_TEXT
            box.line.color.rgb = QUALCOMM_BLUE
            box.line.width = Pt(1.5)

            # Step number + name
            btf = box.text_frame
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            bp = btf.paragraphs[0]
            bp.text = f"{step_num}. {step_name}"
            bp.font.size = Pt(13)
            bp.font.bold = True
            bp.font.color.rgb = text_color
            bp.font.name = BODY_FONT
            bp.alignment = PP_ALIGN.CENTER

            # Arrow between boxes (except last in row)
            if i < 4:
                arrow_x = x + box_w
                arrow_y = y + box_h / 2 - Inches(0.1)
                arrow_box = slide.shapes.add_textbox(arrow_x, arrow_y, gap, Inches(0.3))
                atf = arrow_box.text_frame
                ap = atf.paragraphs[0]
                ap.text = "\u2192"
                ap.font.size = Pt(18)
                ap.font.color.rgb = QUALCOMM_BLUE
                ap.font.name = BODY_FONT
                ap.alignment = PP_ALIGN.CENTER

    # Connecting arrow between rows (right side of row 1 to left of row 2)
    conn_x = start_x + 4 * (box_w + gap) + box_w + Inches(0.1)
    conn_y = y1 + box_h
    conn_box = slide.shapes.add_textbox(conn_x - Inches(0.5), conn_y, Inches(1.0), Inches(1.8))
    ctf = conn_box.text_frame
    cp = ctf.paragraphs[0]
    cp.text = "\u21B5"
    cp.font.size = Pt(28)
    cp.font.color.rgb = QUALCOMM_BLUE
    cp.font.name = BODY_FONT
    cp.alignment = PP_ALIGN.CENTER

    # Command line
    cmd_top = Inches(5.8)
    cmd_box = slide.shapes.add_textbox(Inches(2.0), cmd_top, Inches(9.0), Inches(0.5))
    tf = cmd_box.text_frame
    p = tf.paragraphs[0]
    p.text = "pip install qualcomm-ai-toolkit && quad quickstart"
    p.font.size = Pt(16)
    p.font.name = "Consolas"
    p.font.color.rgb = QUALCOMM_BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tag_top = Inches(6.4)
    tag_box = slide.shapes.add_textbox(Inches(2.0), tag_top, Inches(9.0), Inches(0.5))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "From discovery to production in one platform."
    p.font.size = Pt(18)
    p.font.name = BODY_FONT
    p.font.color.rgb = DARK_TEXT
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 12)


if __name__ == "__main__":
    create_presentation()
